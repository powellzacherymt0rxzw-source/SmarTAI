from io import BytesIO

import fitz
import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from fastapi.testclient import TestClient


def _docx_bytes() -> bytes:
    document = Document()
    document.add_paragraph("DOCX knowledge content")
    stream = BytesIO()
    document.save(stream)
    return stream.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    box.text = "PPTX knowledge content"
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _pdf_bytes() -> bytes:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "PDF knowledge content")
    content = document.tobytes()
    document.close()
    return content


@pytest.mark.asyncio
@pytest.mark.parametrize(
    "filename,body,expected",
    [
        ("notes.txt", b"TXT knowledge content", "TXT knowledge content"),
        ("notes.md", b"# Markdown knowledge content", "Markdown knowledge content"),
        ("notes.pdf", _pdf_bytes(), "PDF knowledge content"),
        ("notes.docx", _docx_bytes(), "DOCX knowledge content"),
        ("notes.pptx", _pptx_bytes(), "PPTX knowledge content"),
    ],
    ids=["txt", "md", "pdf", "docx", "pptx"],
)
async def test_extract_text_supports_personal_knowledge_formats(filename, body, expected):
    from backend.rag.chunker import extract_text

    text = await extract_text(filename, body)
    assert expected in text


@pytest.mark.asyncio
async def test_extract_text_rejects_unsupported_personal_knowledge_format():
    from fastapi import HTTPException
    from backend.rag.chunker import extract_text

    with pytest.raises(HTTPException) as exc:
        await extract_text("notes.png", b"not supported")
    assert exc.value.status_code == 400


def test_knowledge_repository_is_owner_scoped_and_persists_assignment_selection():
    from backend.db.knowledge_repository import (
        create_document,
        get_document,
        list_documents,
        list_selected_documents,
        replace_document_chunks,
        set_task_documents,
    )
    from backend.db.models import AssignmentRecord, CourseRecord, UserRecord
    from backend.db.session import session_scope

    with session_scope() as session:
        session.add_all([
            UserRecord(id="owner-a", username="owner-a", password_hash="hash", role="teacher",
                       is_active=True),
            UserRecord(id="owner-b", username="owner-b", password_hash="hash", role="teacher",
                       is_active=True),
        ])
        session.flush()
        # A course + assignment owned by owner-a; selection is assignment-scoped.
        session.add(CourseRecord(id="course-a", name="C", teacher_id="owner-a"))
        session.flush()
        session.add(AssignmentRecord(id="asg-a", course_id="course-a", teacher_id="owner-a",
                                     name="A", status="draft", version=1))

    document = create_document(
        document_id="doc-a", owner_id="owner-a", stored_file_id=None,
        title="Algebra notes", original_name="notes.txt", content_type="text/plain",
        size_bytes=12, sha256="a" * 64,
    )
    replace_document_chunks(document.id, ["quadratic formula", "factorization"])

    assert get_document("doc-a", "owner-a").status == "ready"
    assert get_document("doc-a", "owner-b") is None
    assert [item.id for item in list_documents("owner-a")] == ["doc-a"]
    assert list_documents("owner-b") == []

    set_task_documents(assignment_id="asg-a", owner_id="owner-a", document_ids=["doc-a"])
    assert [item.id for item in list_selected_documents("asg-a", "owner-a")] == ["doc-a"]
    # Cross-owner selection resolves nothing (assignment owner predicate).
    assert list_selected_documents("asg-a", "owner-b") == []


@pytest.mark.asyncio
async def test_persistent_bm25_retriever_only_uses_selected_documents():
    from backend.db.knowledge_repository import create_document, replace_document_chunks, set_task_documents
    from backend.db.models import AssignmentRecord, CourseRecord, UserRecord
    from backend.db.session import session_scope
    from backend.knowledge.retriever import PersistentKnowledgeRetriever

    with session_scope() as session:
        session.add(UserRecord(id="retriever-owner", username="retriever-owner", password_hash="hash",
                               role="teacher", is_active=True))
        session.flush()
        session.add(CourseRecord(id="rc", name="C", teacher_id="retriever-owner"))
        session.flush()
        session.add(AssignmentRecord(id="retriever-asg", course_id="rc", teacher_id="retriever-owner",
                                     name="A", status="draft", version=1))

    selected = create_document(document_id="selected", owner_id="retriever-owner", stored_file_id=None,
                               title="Selected", original_name="selected.txt", content_type="text/plain",
                               size_bytes=1, sha256="b" * 64)
    hidden = create_document(document_id="hidden", owner_id="retriever-owner", stored_file_id=None,
                             title="Hidden", original_name="hidden.txt", content_type="text/plain",
                             size_bytes=1, sha256="c" * 64)
    replace_document_chunks(selected.id, ["Newton mechanics acceleration force"])
    replace_document_chunks(hidden.id, ["Newton secret unselected sentence"])
    set_task_documents(assignment_id="retriever-asg", owner_id="retriever-owner", document_ids=[selected.id])

    # Scope is now the assignment id; the retriever resolves the teacher owner
    # through the assignment and only ranks the selected document's chunks.
    results = await PersistentKnowledgeRetriever().retrieve("Newton force", k=5, scope="retriever-asg")

    assert results
    assert results[0].source == "selected.txt"
    assert all("secret" not in item.content for item in results)


def test_personal_knowledge_api_upload_list_download_and_delete():
    from backend.auth import create_token
    from backend.db.models import UserRecord
    from backend.db.session import session_scope
    from backend.main import app

    with session_scope() as session:
        session.add(UserRecord(id="demo_api-kb-owner", username="api-kb-owner", password_hash="hash",
                               role="teacher", is_active=True))

    client = TestClient(app)
    headers = {"Authorization": f"Bearer {create_token('demo_api-kb-owner', 'teacher')}"}
    uploaded = client.post("/knowledge/documents", headers=headers,
                           files={"file": ("api-notes.txt", b"persistent API knowledge", "text/plain")})
    assert uploaded.status_code == 201, uploaded.text
    document = uploaded.json()
    assert document["status"] == "ready"

    listed = client.get("/knowledge/documents", headers=headers)
    assert listed.status_code == 200
    assert [item["id"] for item in listed.json()["documents"]] == [document["id"]]

    downloaded = client.get(f"/knowledge/documents/{document['id']}/download", headers=headers)
    assert downloaded.status_code == 200
    assert downloaded.content == b"persistent API knowledge"

    assert client.get(f"/knowledge/documents/{document['id']}").status_code == 401
    removed = client.delete(f"/knowledge/documents/{document['id']}", headers=headers)
    assert removed.status_code == 200
    assert client.get(f"/knowledge/documents/{document['id']}", headers=headers).status_code == 404


def test_assignment_knowledge_selection_caps_at_three_ready_documents():
    """A teacher may select at most three ready personal documents per assignment,
    and the selection survives repository recreation (DB is the source of truth)."""
    from backend.db.knowledge_repository import (
        create_document, list_selected_documents, replace_document_chunks, set_task_documents,
    )
    from backend.db.models import AssignmentRecord, CourseRecord, UserRecord
    from backend.db.session import session_scope

    with session_scope() as session:
        session.add(UserRecord(id="cap-owner", username="cap-owner", password_hash="hash", role="teacher", is_active=True))
        session.flush()
        session.add(CourseRecord(id="cap-c", name="C", teacher_id="cap-owner"))
        session.flush()
        session.add(AssignmentRecord(id="cap-asg", course_id="cap-c", teacher_id="cap-owner", name="A",
                                     status="draft", version=1))

    doc_ids = []
    for i in range(4):
        doc = create_document(document_id=f"cap-doc-{i}", owner_id="cap-owner", stored_file_id=None,
                              title=f"D{i}", original_name=f"d{i}.txt", content_type="text/plain",
                              size_bytes=1, sha256=hex(i * 17 + 1)[2:].rjust(64, "0"))
        replace_document_chunks(doc.id, [f"chunk content {i}"])
        doc_ids.append(doc.id)

    # Selecting four must be rejected.
    with pytest.raises(ValueError):
        set_task_documents(assignment_id="cap-asg", owner_id="cap-owner", document_ids=doc_ids)

    # Three is the cap and persists across a fresh repository call.
    set_task_documents(assignment_id="cap-asg", owner_id="cap-owner", document_ids=doc_ids[:3])
    assert [d.id for d in list_selected_documents("cap-asg", "cap-owner")] == doc_ids[:3]
